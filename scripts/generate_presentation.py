import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    bg_color = RGBColor(15, 23, 42)        # #0f172a (Dark Navy)
    title_color = RGBColor(255, 255, 255)   # #ffffff (White)
    accent_color = RGBColor(192, 132, 252) # #c084fc (Light Purple)
    text_color = RGBColor(226, 232, 240)   # #e2e8f0 (Light Gray)
    
    def add_blank_slide():
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        return slide

    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = title_color
        return title_box

    def add_bullets(slide, items, left, top, width, height, font_size=14):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(8)
            p.line_spacing = 1.2
            p.font.name = 'Segoe UI'
            p.font.size = Pt(font_size)
            p.font.color.rgb = text_color
        return txBox

    # ----------------------------------------------------
    # Slide 1: Cover
    # ----------------------------------------------------
    slide = add_blank_slide()
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    p1.text = "Service Mesh & Observability\nfor Microservices"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = title_color
    p1.space_after = Pt(18)
    
    p2 = tf.add_paragraph()
    p2.text = "MeshMart: Resiliency, Traffic Control, and Telemetry in a Kubernetes Cluster"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(18)
    p2.font.color.rgb = accent_color
    p2.space_after = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "Distributed Systems Course Project | Presentation Deck"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = 'Segoe UI'
    p3.font.size = Pt(12)
    p3.font.color.rgb = text_color

    # ----------------------------------------------------
    # Slide 2: Problem Statement
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "The Monolith to Microservices Paradox")
    bullets_s2 = [
        "- Dynamic Business Split: Monolith decomposition creates isolated scaling domains but introduces network dependencies.",
        "- Network Boundary Failures: Internal function calls are replaced by HTTP network hops that can lag, drop, or fail.",
        "- Operational Blindspots: Difficult to trace request latency or isolate failure locations across many backend microservices.",
        "- Code Clutter: Implementing retry, timeout, and circuit breaker logic in business code leads to duplication and complexity.",
        "- Infrastructure Solution: Moving communication policies to the platform layer keeps service microservices clean."
    ]
    add_bullets(slide, bullets_s2, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 3: Scope & Service Boundaries
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "MeshMart Scope & Service Boundaries")
    bullets_s3 = [
        "- frontend: Nginx container serving the storefront web UI and proxying endpoints.",
        "- product-service: FastAPI app managing database catalog, reviews, and stock locks.",
        "- order-service: FastAPI app acting as the transactional Saga orchestrator and transaction logger.",
        "- payment-service: FastAPI app simulating bank approvals, declines, and processor delays.",
        "- notification-service: FastAPI app simulating SMS or email alerts to consumers post-checkout."
    ]
    add_bullets(slide, bullets_s3, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 4: API Endpoints & Contracts
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "API Endpoints & Service Contracts")
    bullets_s4 = [
        "- GET /products ➔ Returns catalog products list.",
        "- POST /products/{id}/reviews ➔ Submits customer review details.",
        "- POST /orders ➔ Initiates order checklist flow (supports Idempotency-Key headers).",
        "- POST /inventory/reserve ➔ Locks product inventory in database.",
        "- POST /inventory/commit ➔ Confirms permanent inventory decrement on payment success.",
        "- POST /payments ➔ Simulates transaction billing.",
        "- POST /notifications ➔ Sends alert result directly."
    ]
    add_bullets(slide, bullets_s4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=14)

    # ----------------------------------------------------
    # Slide 5: System Deployment Architecture (With Diagram)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "System Deployment Architecture")
    bullets_s5 = [
        "- Ingress Gateway: Public traffic enters via the Istio Ingress Gateway and maps routes by URL path.",
        "- Envoy Interception: Each application container is paired with an Envoy sidecar proxy (istio-proxy) injected in the meshmart namespace.",
        "- Control Plane (istiod): Translates high-level YAML policies into xDS config updates for proxies.",
        "- Observability Pipeline: Sidecar Envoy proxies automatically collect and emit performance telemetry to metrics stores."
    ]
    add_bullets(slide, bullets_s5, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), font_size=14)
    
    img_path_s5 = 'docs/system_architecture.png'
    if os.path.exists(img_path_s5):
        slide.shapes.add_picture(img_path_s5, Inches(6.8), Inches(1.8), width=Inches(5.733), height=Inches(4.5))

    # ----------------------------------------------------
    # Slide 6: How Envoy Sidecars Work (With Diagram)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Network Interception: How Envoy Works")
    bullets_s6 = [
        "- iptables Redirection: All traffic entering or leaving a pod is transparently redirected to local Envoy proxies.",
        "- Inbound Traffic: Ingress calls are intercepted by Envoy, checked for security policies, and forwarded via localhost.",
        "- Outbound Traffic: App outbound calls are intercepted by Envoy to apply timeouts, retries, and trace headers.",
        "- Zero App Intrusion: FastAPI applications operate without modifying code to interact with network rules."
    ]
    add_bullets(slide, bullets_s6, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), font_size=14)
    
    img_path_s6 = 'docs/envoy_sidecar_work.png'
    if os.path.exists(img_path_s6):
        slide.shapes.add_picture(img_path_s6, Inches(6.8), Inches(1.8), width=Inches(5.733), height=Inches(4.5))

    # ----------------------------------------------------
    # Slide 7: Saga Pattern (Distributed Transactions)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Distributed Transactions (Saga Pattern)")
    bullets_s7 = [
        "- Isolation Limits: Traditional database locks cannot span network boundaries safely in microservices.",
        "- Saga Orchestration: order-service manages the checkout steps sequentially.",
        "- Eventual Consistency: Each microservice processes its local transaction and reports state back.",
        "- Dynamic Rollback: If payment declines, the orchestrator triggers compensation calls to rollback stock."
    ]
    add_bullets(slide, bullets_s7, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 8: Step-by-Step Saga Success Path (With Diagram)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Saga Transaction Success Flow")
    bullets_s8 = [
        "- Check Idempotency: Checks database key to prevent double checkout processing.",
        "- Step 1 (Reserve): order-service requests /inventory/reserve to lock catalog stock.",
        "- Step 2 (Billing): order-service queries payment-service (/payments).",
        "- Step 3 (Commit): On payment success, order-service calls /inventory/commit to complete stock reduction.",
        "- Step 4 (Alert): order-service notifies user via notification-service."
    ]
    add_bullets(slide, bullets_s8, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), font_size=14)
    
    img_path_s8 = 'docs/saga_checkout_flow.png'
    if os.path.exists(img_path_s8):
        slide.shapes.add_picture(img_path_s8, Inches(6.8), Inches(1.8), width=Inches(5.733), height=Inches(4.5))

    # ----------------------------------------------------
    # Slide 9: Saga Compensation & Rollback (With Diagram)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Saga Rollback & Stock Compensation")
    bullets_s9 = [
        "- Problem: Stock lock succeeded, but credit card transaction failed or timed out.",
        "- Inconsistent State: Products cannot remain locked if payment failed.",
        "- Step 1 (Release): order-service catches payment failure and calls /inventory/release.",
        "- Step 2 (Restore): product-service restores in-memory stock values.",
        "- Result: Avoids database drift or stock locking issues during billing failures."
    ]
    add_bullets(slide, bullets_s9, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), font_size=14)
    
    if os.path.exists(img_path_s8):
        slide.shapes.add_picture(img_path_s8, Inches(6.8), Inches(1.8), width=Inches(5.733), height=Inches(4.5))

    # ----------------------------------------------------
    # Slide 10: Idempotency Guard (Preventing Double Billing)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Idempotency Key Guard")
    bullets_s10 = [
        "- Network Disconnection Risk: A client billing request succeeds, but connection drops before response is received.",
        "- The Double Click Problem: Users click 'Checkout' twice, executing duplicate payments.",
        "- Idempotency Cache: order-service caches Idempotency-Key headers with responses.",
        "- Fingerprint Verification: Verifies order contents match the original payload.",
        "- Conflict handling: Returns active processing errors (409) or cached order logs immediately."
    ]
    add_bullets(slide, bullets_s10, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 11: Istio Ingress Gateway & Routing
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Ingress Path-Based Routing")
    bullets_s11 = [
        "- Gateway Configuration: Single public host/port receives external client calls.",
        "- Path Prefixes (virtual-service): Matches endpoints to backend services.",
        "- Route mappings: / ➔ frontend, /products ➔ product-service, /orders ➔ order-service.",
        "- Decoupled Architecture: Public URLs are clean and decoupled from target cluster ports."
    ]
    add_bullets(slide, bullets_s11, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 12: Traffic Control: Timeouts & Retries
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Resilience: Timeouts & Retries")
    bullets_s12 = [
        "- Timeouts: Aborts lagging calls at proxy levels (orders: 8s, product: 4s, payments: 3s).",
        "- Protects thread pools: Prevent slow downstream dependencies from locking server queues.",
        "- Retries: 2 retry attempts configured in VirtualServices.",
        "- Transient Failures: Envoy transparently handles socket glitches before reporting errors to app."
    ]
    add_bullets(slide, bullets_s12, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 13: Fault Injection (Chaos Engineering)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Chaos Engineering: Fault Injection")
    bullets_s13 = [
        "- Testing resilience: Injecting failures into live staging environments without changing code.",
        "- Istio Fault policy: Configured in payment-fault-injection.yaml.",
        "- Settings: 30% of requests to payment-service receive a 2-second delay.",
        "- Validation: Confirms order-service handles timeout limits gracefully."
    ]
    add_bullets(slide, bullets_s13, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 14: DestinationRules: Connection Pools
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "DestinationRules: Connection Pools")
    bullets_s14 = [
        "- Traffic Policy Settings: Caps max TCP/HTTP connection resources in destination-rule.yaml.",
        "- Port Protection: Prevents microservices from crash failures during high traffic surges.",
        "- Configured Limits: payment-service: 20 max connections, product-service: 50, order-service: 100.",
        "- Client queueing: Envoy queues overflow requests until downstream resources clear."
    ]
    add_bullets(slide, bullets_s14, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 15: Outlier Detection (Circuit Breaking)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Outlier Detection (Circuit Breaking)")
    bullets_s15 = [
        "- Self-Healing Mesh: Monitors and ejects bad replica pods automatically.",
        "- Trigger settings: 3 consecutive 5xx errors in a 10s interval.",
        "- Ejection Rule: Bypasses the bad pod for 30 seconds, forwarding traffic to healthy replicas.",
        "- Cascade Protection: Stops failing nodes from crashing other dependencies."
    ]
    add_bullets(slide, bullets_s15, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 16: Observability Data Flow (With Diagram)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Observability & Telemetry Pipeline")
    bullets_s16 = [
        "- Prometheus Scraper: Gathers HTTP request metrics directly from sidecars every 15s.",
        "- Jaeger Collector: Receives span spans to trace requests across microservice hops.",
        "- Kiali Graph: Consumes metrics and configs to draw a live dependency graph.",
        "- Grafana: Queries metrics and displays dashboards (Latency percentiles, Traffic, Error Rates)."
    ]
    add_bullets(slide, bullets_s16, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), font_size=14)
    
    img_path_s16 = 'docs/observability_data_flow.png'
    if os.path.exists(img_path_s16):
        slide.shapes.add_picture(img_path_s16, Inches(6.8), Inches(1.8), width=Inches(5.733), height=Inches(4.5))

    # ----------------------------------------------------
    # Slide 17: Observability in Action: Diagnostics
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Outage Diagnostics & Evidence")
    bullets_s17 = [
        "- Step 1 (Alert): Grafana dashboards show a spike in p90 response latencies.",
        "- Step 2 (Isolate): Kiali mesh graph highlights the failing service edge in red.",
        "- Step 3 (Pinpoint): Jaeger tracing isolates the exact slow/failed downstream span.",
        "- Step 4 (Inspect): Query application logs with the request_id to find trace details.",
        "- Evidence-based: Diagnostics are backed by actual telemetry graphs."
    ]
    add_bullets(slide, bullets_s17, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 18: Scaling & Resource Isolation (With Diagram)
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Workload Scaling & Capacity")
    bullets_s18 = [
        "- Resource Profiles: Pod configurations enforce CPU/Memory limits to prevent starvations.",
        "- Horizontal Pod Autoscaler (HPA): Spawns replica pods dynamically based on CPU request loads.",
        "- Pod Disruption Budget (PDB): Guarantees at least 1 active pod during cluster maintenance.",
        "- Stress Testing: Ran k6 load scripts to trace HPA performance under load."
    ]
    add_bullets(slide, bullets_s18, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), font_size=14)
    
    img_path_s18 = 'docs/traffic_management.png'
    if os.path.exists(img_path_s18):
        slide.shapes.add_picture(img_path_s18, Inches(6.8), Inches(1.8), width=Inches(5.733), height=Inches(4.5))

    # ----------------------------------------------------
    # Slide 19: Limitations & Production Roadmap
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Limitations & Production Roadmap")
    bullets_s19 = [
        "- Persistent Database Integration: Migrate current in-memory lists to PostgreSQL instances.",
        "- Shared Idempotency Cache: Migrate local dict stores to a shared Redis cluster.",
        "- Zero-Trust Security: Enable strict mutual TLS (mTLS) policies and authorization header tokens.",
        "- Automated Alerting: Define Prometheus alert rules with PagerDuty webhook destinations."
    ]
    add_bullets(slide, bullets_s19, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # ----------------------------------------------------
    # Slide 20: Conclusion
    # ----------------------------------------------------
    slide = add_blank_slide()
    add_title(slide, "Key Technical Takeaways")
    bullets_s20 = [
        "- Separation of Concerns: Communication policies are managed by Istio, keeping code lightweight.",
        "- Actionable Telemetry: Combining metrics, traces, and logs is mandatory for operating microservices.",
        "- Resilient Transactions: Distributed systems must adopt Saga rollbacks and idempotency guards.",
        "- Operation Mindset: MeshMart demonstrates how modern services are observed, secured, and scaled."
    ]
    add_bullets(slide, bullets_s20, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), font_size=15)

    # Save presentation
    output_path = 'docs/MeshMart_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
